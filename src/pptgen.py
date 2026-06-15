"""PowerPoint deck generator for the Presales AI Platform.

Builds a 15-slide .pptx from RFP analysis and selected products.
Each AI slide makes exactly one Ollama call to granite4.1:30b and
parses the response via XML markers — same pattern as proposal.py.
"""

import json
import logging
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any

import requests
from dotenv import load_dotenv
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

load_dotenv()

logger = logging.getLogger(__name__)

# ── Model defaults ─────────────────────────────────────────────────────────────

_PPT_MODEL = os.getenv("OLLAMA_COMPLIANCE_MODEL", "granite4.1:30b")
_OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")

# ── Brand colours ──────────────────────────────────────────────────────────────

_NAVY = RGBColor(0x1F, 0x38, 0x64)       # #1F3864 — header backgrounds
_TEAL = RGBColor(0x00, 0x97, 0xA7)       # #0097A7 — accent bar, table headers
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7) # alternate table row
_BODY_TEXT = RGBColor(0x1A, 0x1A, 0x2E)  # near-black body copy
_GREY = RGBColor(0x60, 0x60, 0x60)       # footer / caption text

_NAVY_HEX = "1F3864"
_TEAL_HEX = "0097A7"
_LIGHT_GREY_HEX = "F2F4F7"

# ── Slide geometry (16:9, 10" × 5.625") ────────────────────────────────────────

_SLIDE_W = Inches(10)
_SLIDE_H = Inches(5.625)

# Title text box — sits below the teal accent bar
_TITLE_TOP = Inches(0.45)
_TITLE_LEFT = Inches(0.4)
_TITLE_WIDTH = Inches(9.2)
_TITLE_HEIGHT = Inches(0.75)

# Body area — below title
_BODY_TOP = Inches(1.35)
_BODY_LEFT = Inches(0.4)
_BODY_WIDTH = Inches(9.2)
_BODY_HEIGHT = Inches(3.9)

# Teal accent bar — 3px strip at top of slide
_BAR_TOP = Emu(0)
_BAR_LEFT = Emu(0)
_BAR_WIDTH = _SLIDE_W
_BAR_HEIGHT = Inches(0.06)

# ── System prompt ──────────────────────────────────────────────────────────────

_SYS_PPT = (
    "You are a professional presales presentation writer for Nexus Group, "
    "specialising in PKI, digital identity, and cybersecurity. "
    "Write concise, punchy slide content — short bullet points, no waffle. "
    "Always wrap your response in the specified XML tags and nothing else."
)


# ── Ollama and XML helpers ─────────────────────────────────────────────────────

def _call_ollama(
    model: str,
    prompt: str,
    base_url: str,
    system: str = "",
    timeout: int = 180,
    num_predict: int = 512,
) -> str:
    """Call the Ollama generate API and return the raw response text.

    Args:
        model: Ollama model name.
        prompt: User prompt text.
        base_url: Ollama API base URL.
        system: Optional system prompt string.
        timeout: HTTP timeout in seconds.
        num_predict: Max tokens to generate.

    Returns:
        Raw model response text, stripped of whitespace.

    Raises:
        requests.exceptions.RequestException: On network or HTTP errors.
    """
    payload: dict[str, Any] = {
        "model": model,
        "prompt": prompt,
        "stream": False,
        "options": {"temperature": 0.3, "num_predict": num_predict},
    }
    if system:
        payload["system"] = system
    logger.info("Calling Ollama '%s' for PPT slide", model)
    response = requests.post(
        f"{base_url}/api/generate", json=payload, timeout=timeout
    )
    response.raise_for_status()
    return response.json().get("response", "").strip()


def _extract_xml(text: str, tag: str) -> str:
    """Extract content between <tag> and </tag>; falls back to full text.

    Args:
        text: Raw model response.
        tag: XML tag name without angle brackets.

    Returns:
        Extracted content, stripped.
    """
    match = re.search(f"<{tag}>(.*?)</{tag}>", text, re.DOTALL)
    if match:
        return match.group(1).strip()
    logger.warning("XML tag <%s> not found in model response; using raw text", tag)
    return text.strip()


def _slim_analysis(analysis: dict[str, Any]) -> dict[str, Any]:
    """Return a trimmed copy of the analysis dict to keep prompts compact.

    Caps executive_summary at 250 chars, requirements at top 8 with each
    trimmed to 120 chars, win_themes and risk_areas at top 3.

    Args:
        analysis: Full RFP analysis dict.

    Returns:
        Shallow-copied dict with trimmed fields.
    """
    slim = dict(analysis)
    slim["executive_summary"] = analysis.get("executive_summary", "")[:250]
    slim["win_themes"] = analysis.get("win_themes", [])[:3]
    slim["risk_areas"] = analysis.get("risk_areas", [])[:3]
    slim_reqs = []
    for r in analysis.get("requirements", [])[:8]:
        slim_r = dict(r)
        slim_r["requirement"] = r.get("requirement", "")[:120]
        slim_reqs.append(slim_r)
    slim["requirements"] = slim_reqs
    return slim


# ── python-pptx layout helpers ─────────────────────────────────────────────────

def _new_blank_slide(prs: Presentation) -> Any:
    """Add a blank slide (layout index 6) and return it.

    Args:
        prs: The Presentation to add a slide to.

    Returns:
        The new Slide object.
    """
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _fill_shape(shape: Any, rgb: RGBColor) -> None:
    """Set a shape's solid fill colour.

    Args:
        shape: python-pptx shape with a fill attribute.
        rgb: RGBColor to apply.
    """
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_teal_bar(slide: Any) -> None:
    """Add the 6pt teal accent stripe across the very top of the slide.

    Args:
        slide: Target Slide.
    """
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        _BAR_LEFT, _BAR_TOP, _BAR_WIDTH, _BAR_HEIGHT,
    )
    _fill_shape(bar, _TEAL)
    bar.line.fill.background()  # no border


def _set_run_font(run: Any, size_pt: int, bold: bool, colour: RGBColor) -> None:
    """Apply font size, bold, and colour to a text run.

    Args:
        run: python-pptx text run.
        size_pt: Font size in points.
        bold: Whether to bold the run.
        colour: RGBColor for the run.
    """
    run.font.name = "Calibri"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = colour


def set_slide_title(slide: Any, text: str) -> None:
    """Add a white Calibri Bold 28pt title text box below the teal bar.

    Args:
        slide: Target Slide.
        text: Title text.
    """
    txb = slide.shapes.add_textbox(
        _TITLE_LEFT, _TITLE_TOP, _TITLE_WIDTH, _TITLE_HEIGHT
    )
    tf = txb.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    _set_run_font(run, 28, True, _NAVY)


def add_bullet_slide(
    slide: Any,
    bullets: list[str],
    font_size: int = 18,
    top: Emu | None = None,
    left: Emu | None = None,
    width: Emu | None = None,
    height: Emu | None = None,
) -> None:
    """Add a bullet-point text box to the slide body area.

    Args:
        slide: Target Slide.
        bullets: List of bullet strings.
        font_size: Body font size in points.
        top / left / width / height: Override geometry; defaults to _BODY_* constants.
    """
    t = top if top is not None else _BODY_TOP
    l = left if left is not None else _BODY_LEFT
    w = width if width is not None else _BODY_WIDTH
    h = height if height is not None else _BODY_HEIGHT

    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = Pt(4)
        para.level = 0
        run = para.add_run()
        run.text = f"• {bullet}"
        _set_run_font(run, font_size, False, _BODY_TEXT)


def add_table_slide(
    slide: Any,
    headers: list[str],
    rows: list[list[str]],
    top: Emu | None = None,
    left: Emu | None = None,
    width: Emu | None = None,
    height: Emu | None = None,
) -> None:
    """Add a branded table (teal header, alternating rows) to a slide.

    Args:
        slide: Target Slide.
        headers: Column header labels.
        rows: Data rows (list of string lists).
        top / left / width / height: Override geometry; defaults to _BODY_* constants.
    """
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    t = top if top is not None else _BODY_TOP
    l = left if left is not None else _BODY_LEFT
    w = width if width is not None else _BODY_WIDTH
    h = height if height is not None else _BODY_HEIGHT

    n_cols = len(headers)
    n_rows = 1 + len(rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols, l, t, w, h)
    tbl = table_shape.table

    # Header row — teal background, white bold text
    for col_idx, header in enumerate(headers):
        cell = tbl.cell(0, col_idx)
        cell.text = header
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.text = header
        _set_run_font(run, 11, True, _WHITE)
        # Teal fill via XML
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", _TEAL_HEX)

    # Data rows
    for row_idx, row_vals in enumerate(rows):
        fill_hex = _LIGHT_GREY_HEX if row_idx % 2 == 0 else "FFFFFF"
        for col_idx, val in enumerate(row_vals):
            cell = tbl.cell(row_idx + 1, col_idx)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            run.text = str(val)
            _set_run_font(run, 10, False, _BODY_TEXT)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
            srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
            srgbClr.set("val", fill_hex)


def add_section_header(prs: Presentation, title: str, subtitle: str = "") -> None:
    """Add a full-bleed navy section-divider slide.

    Used to visually separate major blocks in the deck. White centred text
    on a navy background with a teal accent bar.

    Args:
        prs: The Presentation.
        title: Large white centred title text.
        subtitle: Optional smaller subtitle below the title.
    """
    slide = _new_blank_slide(prs)

    # Navy background
    bg = slide.shapes.add_shape(1, Emu(0), Emu(0), _SLIDE_W, _SLIDE_H)
    _fill_shape(bg, _NAVY)
    bg.line.fill.background()

    # Teal bar
    _add_teal_bar(slide)

    # Centred title
    txb = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8), Inches(9), Inches(1.2)
    )
    tf = txb.text_frame
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = title
    _set_run_font(run, 36, True, _WHITE)

    if subtitle:
        sub_para = tf.add_paragraph()
        sub_para.alignment = PP_ALIGN.CENTER
        sub_run = sub_para.add_run()
        sub_run.text = subtitle
        _set_run_font(sub_run, 18, False, RGBColor(0xCC, 0xE5, 0xFF))


# ── Slide generators: slides 1–7 ──────────────────────────────────────────────

def slide_01_title(
    prs: Presentation,
    analysis: dict[str, Any],
    selected_products: list[dict[str, Any]],
) -> None:
    """Slide 1: Cover slide. Static — no Ollama call.

    Full-bleed navy background. Customer name, Nexus branding, date,
    product list, and Confidential tag.

    Args:
        prs: The Presentation.
        analysis: RFP analysis dict (uses _rfp_filename).
        selected_products: Selected product dicts.
    """
    slide = _new_blank_slide(prs)

    # Navy background
    bg = slide.shapes.add_shape(1, Emu(0), Emu(0), _SLIDE_W, _SLIDE_H)
    _fill_shape(bg, _NAVY)
    bg.line.fill.background()

    # Teal accent bar at top
    _add_teal_bar(slide)

    rfp_filename = analysis.get("_rfp_filename", "rfp")
    customer_name = Path(rfp_filename).stem.replace("_", " ").replace("-", " ").title()
    product_names = (
        "  |  ".join(p.get("name", "") for p in selected_products)
        if selected_products
        else "Nexus PKI Solution"
    )
    today_str = datetime.now().strftime("%d %B %Y")

    # "NEXUS GROUP" brand label
    brand_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.7), Inches(9.2), Inches(0.4))
    para = brand_box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = "NEXUS GROUP"
    _set_run_font(run, 11, True, _TEAL)

    # Main title — customer name
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.4), Inches(9.2), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = f"Technical Proposal — {customer_name}"
    _set_run_font(run, 32, True, _WHITE)

    # Product subtitle
    prod_box = slide.shapes.add_textbox(Inches(0.4), Inches(2.8), Inches(9.2), Inches(0.6))
    para = prod_box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = product_names
    _set_run_font(run, 14, False, RGBColor(0xCC, 0xE5, 0xFF))

    # Meta info (date, version, classification)
    meta_lines = [
        f"Date:  {today_str}",
        "Version:  1.0 DRAFT",
        "Classification:  Confidential",
    ]
    meta_box = slide.shapes.add_textbox(Inches(0.4), Inches(4.3), Inches(9.2), Inches(1.0))
    tf = meta_box.text_frame
    for i, line in enumerate(meta_lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        _set_run_font(run, 11, False, _GREY)


def slide_02_agenda(prs: Presentation) -> None:
    """Slide 2: Agenda. Static — no Ollama call.

    Two-column numbered list of slides 3–15.

    Args:
        prs: The Presentation.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Agenda")

    left_items = [
        "1.  Executive Summary",
        "2.  Customer Challenge",
        "3.  Proposed Solution Overview",
        "4.  Product Portfolio",
        "5.  Product Deep Dive",
        "6.  How We Meet Your Requirements",
        "7.  Compliance Summary",
    ]
    right_items = [
        "8.  Implementation Approach",
        "9.  Project Timeline",
        "10. Our Team",
        "11. Why Nexus",
        "12. References & Case Studies",
        "13. Next Steps",
    ]

    # Left column
    txb_l = slide.shapes.add_textbox(Inches(0.4), _BODY_TOP, Inches(4.4), _BODY_HEIGHT)
    tf = txb_l.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = Pt(5)
        run = para.add_run()
        run.text = item
        _set_run_font(run, 16, False, _BODY_TEXT)

    # Right column
    txb_r = slide.shapes.add_textbox(Inches(5.2), _BODY_TOP, Inches(4.4), _BODY_HEIGHT)
    tf = txb_r.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = Pt(5)
        run = para.add_run()
        run.text = item
        _set_run_font(run, 16, False, _BODY_TEXT)


def slide_03_exec_summary(
    prs: Presentation,
    analysis: dict[str, Any],
    selected_products: list[dict[str, Any]],
    config: dict[str, Any],
) -> None:
    """Slide 3: Executive Summary. AI-generated via granite4.1:30b.

    Prompts for 4–5 concise bullet points. XML tag: <exec_summary>.

    Args:
        prs: The Presentation.
        analysis: RFP analysis dict.
        selected_products: Selected product dicts.
        config: Runtime config (ollama_base_url, ppt_model).
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Executive Summary")

    base_url = config["ollama_base_url"]
    model = config["ppt_model"]
    slim = _slim_analysis(analysis)
    product_names = ", ".join(p.get("name", "") for p in selected_products)
    industry = analysis.get("industry_vertical", "")

    prompt = (
        "Write exactly 4 concise executive-summary bullet points for a Nexus Group "
        "PKI/digital identity sales presentation. Each bullet must be one short sentence. "
        "Cover: (1) the customer's strategic need, (2) the Nexus solution fit, "
        "(3) a key differentiator, (4) the expected business outcome. "
        "No sub-bullets. No filler words. "
        "Wrap your four bullets in <exec_summary> tags, one bullet per line.\n\n"
        f"Industry: {industry}\n"
        f"Products: {product_names}\n"
        f"RFP summary: {slim['executive_summary']}\n"
        f"Win themes: {json.dumps(slim['win_themes'])}"
    )

    raw = _call_ollama(model, prompt, base_url, system=_SYS_PPT, num_predict=300)
    content = _extract_xml(raw, "exec_summary")
    bullets = [ln.lstrip("•-– ").strip() for ln in content.splitlines() if ln.strip()]
    add_bullet_slide(slide, bullets[:5])


def slide_04_challenge(
    prs: Presentation,
    analysis: dict[str, Any],
    config: dict[str, Any],
) -> None:
    """Slide 4: Customer Challenge. AI-generated via granite4.1:30b.

    Produces 4 bullet points. XML tag: <challenge>.

    Args:
        prs: The Presentation.
        analysis: RFP analysis dict.
        config: Runtime config.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Customer Challenge")

    base_url = config["ollama_base_url"]
    model = config["ppt_model"]
    slim = _slim_analysis(analysis)

    top_reqs = "\n".join(
        f"- {r.get('requirement', '')}"
        for r in slim["requirements"]
        if str(r.get("type", "")).upper() == "MANDATORY"
    )

    prompt = (
        "Write exactly 4 bullet points describing the customer's key business and "
        "technical challenges that this RFP is trying to solve. "
        "Each bullet must be one short, impactful sentence. No sub-bullets. "
        "Wrap in <challenge> tags, one bullet per line.\n\n"
        f"Risk areas: {json.dumps(slim['risk_areas'])}\n"
        f"Top mandatory requirements:\n{top_reqs}"
    )

    raw = _call_ollama(model, prompt, base_url, system=_SYS_PPT, num_predict=250)
    content = _extract_xml(raw, "challenge")
    bullets = [ln.lstrip("•-– ").strip() for ln in content.splitlines() if ln.strip()]
    add_bullet_slide(slide, bullets[:4])


def slide_05_solution_overview(
    prs: Presentation,
    analysis: dict[str, Any],
    selected_products: list[dict[str, Any]],
    config: dict[str, Any],
) -> None:
    """Slide 5: Proposed Solution Overview. AI-generated via granite4.1:30b.

    One intro line + one bullet per selected product. XML tag: <solution_overview>.

    Args:
        prs: The Presentation.
        analysis: RFP analysis dict.
        selected_products: Selected product dicts.
        config: Runtime config.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Proposed Solution Overview")

    base_url = config["ollama_base_url"]
    model = config["ppt_model"]
    slim = _slim_analysis(analysis)
    product_names = ", ".join(p.get("name", "") for p in selected_products)

    prompt = (
        "Write a solution overview for a Nexus Group PKI/digital identity presentation. "
        "First line: one sentence introducing the overall solution (no bullet). "
        f"Then write exactly one bullet point per product, in this order: {product_names}. "
        "Each bullet: 'Product Name — what it does for this customer' (max 12 words). "
        "Wrap the entire response in <solution_overview> tags.\n\n"
        f"Customer objectives: {json.dumps(analysis.get('customer_objectives', [])[:3])}\n"
        f"Win themes: {json.dumps(slim['win_themes'])}"
    )

    raw = _call_ollama(model, prompt, base_url, system=_SYS_PPT, num_predict=350)
    content = _extract_xml(raw, "solution_overview")

    lines = [ln.strip() for ln in content.splitlines() if ln.strip()]
    if not lines:
        lines = [f"Nexus Group proposes {product_names} to meet your requirements."]

    # First line as intro paragraph, rest as bullets
    intro_line = lines[0].lstrip("•-– ")
    bullet_lines = [ln.lstrip("•-– ") for ln in lines[1:]]

    intro_box = slide.shapes.add_textbox(
        _BODY_LEFT, _BODY_TOP, _BODY_WIDTH, Inches(0.5)
    )
    tf = intro_box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = intro_line
    _set_run_font(run, 16, True, _NAVY)

    if bullet_lines:
        add_bullet_slide(
            slide, bullet_lines,
            font_size=16,
            top=Inches(1.95),
        )


def slide_06_product_portfolio(
    prs: Presentation,
    selected_products: list[dict[str, Any]],
    config: dict[str, Any],
) -> None:
    """Slide 6: Product Portfolio table. AI-generated via granite4.1:30b.

    Two-column table: Product Name / One-line value proposition.
    XML tag: <product_portfolio>.

    Args:
        prs: The Presentation.
        selected_products: Selected product dicts.
        config: Runtime config.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Product Portfolio")

    base_url = config["ollama_base_url"]
    model = config["ppt_model"]
    product_names = [p.get("name", "") for p in selected_products]
    products_json = json.dumps(product_names)

    prompt = (
        "For each Nexus product listed, write a one-line value proposition "
        "(max 12 words) suitable for a slide table. "
        "Return one line per product in this exact format:\n"
        "Product Name | value proposition\n"
        "Wrap all lines in <product_portfolio> tags.\n\n"
        f"Products: {products_json}"
    )

    raw = _call_ollama(model, prompt, base_url, system=_SYS_PPT, num_predict=300)
    content = _extract_xml(raw, "product_portfolio")

    rows: list[list[str]] = []
    for line in content.splitlines():
        line = line.strip()
        if "|" in line:
            parts = line.split("|", 1)
            rows.append([parts[0].strip(), parts[1].strip()])
        elif line:
            rows.append([line, ""])

    if not rows:
        rows = [[name, ""] for name in product_names]

    add_table_slide(slide, ["Product", "Value Proposition"], rows)


def slide_07_product_detail(
    prs: Presentation,
    analysis: dict[str, Any],
    selected_products: list[dict[str, Any]],
    config: dict[str, Any],
) -> None:
    """Slide 7: Product Deep Dive. AI-generated via granite4.1:30b.

    3 capability bullets for each of the top 3 selected products.
    XML tag: <product_detail>.

    Args:
        prs: The Presentation.
        analysis: RFP analysis dict.
        selected_products: Selected product dicts (top 3 used).
        config: Runtime config.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Product Deep Dive")

    base_url = config["ollama_base_url"]
    model = config["ppt_model"]
    top3 = selected_products[:3]
    product_names = [p.get("name", "") for p in top3]
    industry = analysis.get("industry_vertical", "")

    prompt = (
        "For each Nexus product listed, write exactly 3 short capability bullet points "
        "(max 10 words each) tailored to the customer's industry. "
        f"Return output in this exact format for each product:\n"
        "PRODUCT: [product name]\n"
        "- bullet 1\n- bullet 2\n- bullet 3\n\n"
        "Wrap the entire response in <product_detail> tags.\n\n"
        f"Industry: {industry}\n"
        f"Products: {json.dumps(product_names)}"
    )

    raw = _call_ollama(model, prompt, base_url, system=_SYS_PPT, num_predict=400)
    content = _extract_xml(raw, "product_detail")

    # Parse PRODUCT: blocks
    blocks = re.split(r"PRODUCT\s*:\s*", content, flags=re.IGNORECASE)
    parsed: list[tuple[str, list[str]]] = []
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        lines = block.splitlines()
        prod_name = lines[0].strip()
        bullets = [
            ln.lstrip("•-– ").strip()
            for ln in lines[1:]
            if ln.strip()
        ]
        parsed.append((prod_name, bullets[:3]))

    # Render: up to 3 products side-by-side in columns
    col_width = Inches(3.0)
    col_gap = Inches(0.1)
    for col_idx, (prod_name, bullets) in enumerate(parsed[:3]):
        col_left = _BODY_LEFT + col_idx * (col_width + col_gap)

        # Column header
        hdr = slide.shapes.add_textbox(col_left, _BODY_TOP, col_width, Inches(0.4))
        para = hdr.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = prod_name
        _set_run_font(run, 13, True, _TEAL)

        # Bullets below header
        add_bullet_slide(
            slide, bullets,
            font_size=14,
            top=Inches(1.85),
            left=col_left,
            width=col_width,
            height=Inches(3.4),
        )


# ── Slide generators: slides 8–15 ─────────────────────────────────────────────

def slide_08_requirements_match(
    prs: Presentation,
    analysis: dict[str, Any],
    selected_products: list[dict[str, Any]],
    config: dict[str, Any],
) -> None:
    """Slide 8: How We Meet Your Requirements. AI-generated via granite4.1:30b.

    Table: Requirement (short) / Nexus Product / How we address it.
    XML tag: <requirements_match>.

    Args:
        prs: The Presentation.
        analysis: RFP analysis dict.
        selected_products: Selected product dicts.
        config: Runtime config.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "How We Meet Your Requirements")

    base_url = config["ollama_base_url"]
    model = config["ppt_model"]
    slim = _slim_analysis(analysis)
    product_names = ", ".join(p.get("name", "") for p in selected_products)

    mandatory = [
        r for r in slim["requirements"]
        if str(r.get("type", "")).upper() == "MANDATORY"
    ][:6]
    reqs_text = "\n".join(
        f"- {r.get('id', 'REQ')}: {r.get('requirement', '')}"
        for r in mandatory
    )

    prompt = (
        "For each requirement listed, write one row in this pipe-delimited format:\n"
        "Requirement (max 8 words) | Nexus Product | How addressed (max 8 words)\n"
        "Return only the rows, no header. "
        "Wrap all rows in <requirements_match> tags.\n\n"
        f"Nexus products available: {product_names}\n"
        f"Mandatory requirements:\n{reqs_text}"
    )

    raw = _call_ollama(model, prompt, base_url, system=_SYS_PPT, num_predict=400)
    content = _extract_xml(raw, "requirements_match")

    rows: list[list[str]] = []
    for line in content.splitlines():
        line = line.strip()
        if "|" in line:
            parts = [p.strip() for p in line.split("|")]
            if len(parts) >= 3:
                rows.append(parts[:3])
            elif len(parts) == 2:
                rows.append([parts[0], parts[1], ""])

    if not rows:
        for req in mandatory:
            rows.append([req.get("requirement", "")[:40], product_names.split(",")[0].strip(), ""])

    add_table_slide(
        slide,
        ["Requirement", "Nexus Product", "How Addressed"],
        rows[:8],
    )


def slide_09_compliance_summary(
    prs: Presentation,
    analysis: dict[str, Any],
    selected_products: list[dict[str, Any]],
    config: dict[str, Any],
) -> None:
    """Slide 9: Compliance Summary. AI-generated via granite4.1:30b.

    Table: Product / Compliance % / Key capabilities claimed.
    XML tag: <compliance_summary>.

    Args:
        prs: The Presentation.
        analysis: RFP analysis dict.
        selected_products: Selected product dicts.
        config: Runtime config.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Compliance Summary")

    base_url = config["ollama_base_url"]
    model = config["ppt_model"]
    slim = _slim_analysis(analysis)
    product_names = [p.get("name", "") for p in selected_products]

    prompt = (
        "For each Nexus product listed, estimate its compliance percentage against "
        "the customer's requirements and list 2 key capabilities it covers. "
        "Return one pipe-delimited row per product:\n"
        "Product Name | Compliance % | Key capability 1; Key capability 2\n"
        "Wrap all rows in <compliance_summary> tags.\n\n"
        f"Products: {json.dumps(product_names)}\n"
        f"Requirements summary: {slim['executive_summary']}"
    )

    raw = _call_ollama(model, prompt, base_url, system=_SYS_PPT, num_predict=350)
    content = _extract_xml(raw, "compliance_summary")

    rows: list[list[str]] = []
    for line in content.splitlines():
        line = line.strip()
        if "|" in line:
            parts = [p.strip() for p in line.split("|")]
            if len(parts) >= 3:
                rows.append(parts[:3])
            elif len(parts) == 2:
                rows.append([parts[0], parts[1], ""])

    if not rows:
        rows = [[name, "TBC", ""] for name in product_names]

    add_table_slide(
        slide,
        ["Product", "Compliance", "Key Capabilities"],
        rows,
    )


def slide_10_implementation(
    prs: Presentation,
    config: dict[str, Any],
) -> None:
    """Slide 10: Implementation Approach. AI-generated via granite4.1:30b.

    Table: Phase / Name / Key activities / Duration.
    XML tag: <implementation>.

    Args:
        prs: The Presentation.
        config: Runtime config.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Implementation Approach")

    base_url = config["ollama_base_url"]
    model = config["ppt_model"]

    prompt = (
        "Write exactly 4 implementation phases for a Nexus Group PKI/digital identity project. "
        "Return one pipe-delimited row per phase:\n"
        "Phase N | Phase Name | Key Activities (max 10 words) | Duration\n"
        "Phases must be: Initiation, Installation, Integration & Testing, Go-Live & Handover. "
        "Wrap all rows in <implementation> tags."
    )

    raw = _call_ollama(model, prompt, base_url, system=_SYS_PPT, num_predict=300)
    content = _extract_xml(raw, "implementation")

    rows: list[list[str]] = []
    for line in content.splitlines():
        line = line.strip()
        if "|" in line:
            parts = [p.strip() for p in line.split("|")]
            if len(parts) >= 4:
                rows.append(parts[:4])

    if not rows:
        rows = [
            ["Phase 1", "Initiation", "Kick-off, planning, environment setup", "2 weeks"],
            ["Phase 2", "Installation", "Software install, HSM config", "3 weeks"],
            ["Phase 3", "Integration & Testing", "API integration, UAT", "4 weeks"],
            ["Phase 4", "Go-Live & Handover", "Production cutover, training", "2 weeks"],
        ]

    add_table_slide(
        slide,
        ["Phase", "Name", "Key Activities", "Duration"],
        rows,
    )


def slide_11_timeline(prs: Presentation) -> None:
    """Slide 11: Project Timeline. Static placeholder — no Ollama call.

    Blank table for SA to fill with actual dates.

    Args:
        prs: The Presentation.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Project Timeline")

    placeholder_rows = [["", "", "", "", ""] for _ in range(6)]
    add_table_slide(
        slide,
        ["Phase", "Activity", "Duration", "Start", "End"],
        placeholder_rows,
    )

    note_box = slide.shapes.add_textbox(
        _BODY_LEFT, Inches(5.1), _BODY_WIDTH, Inches(0.4)
    )
    para = note_box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = "[TO BE COMPLETED — confirm dates with customer]"
    _set_run_font(run, 10, False, _GREY)


def slide_12_our_team(prs: Presentation) -> None:
    """Slide 12: Our Team. Static placeholder — no Ollama call.

    Blank table for account team to populate.

    Args:
        prs: The Presentation.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Our Team")

    placeholder_rows = [["", "", "", ""] for _ in range(5)]
    add_table_slide(
        slide,
        ["Role", "Name", "Expertise", "Availability"],
        placeholder_rows,
    )

    note_box = slide.shapes.add_textbox(
        _BODY_LEFT, Inches(5.1), _BODY_WIDTH, Inches(0.4)
    )
    para = note_box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = "[TO BE COMPLETED — insert delivery team details]"
    _set_run_font(run, 10, False, _GREY)


def slide_13_why_nexus(
    prs: Presentation,
    analysis: dict[str, Any],
    config: dict[str, Any],
) -> None:
    """Slide 13: Why Nexus. AI-generated via granite4.1:30b.

    4–5 differentiator bullet points. XML tag: <why_nexus>.

    Args:
        prs: The Presentation.
        analysis: RFP analysis dict.
        config: Runtime config.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Why Nexus")

    base_url = config["ollama_base_url"]
    model = config["ppt_model"]
    industry = analysis.get("industry_vertical", "")

    prompt = (
        "Write exactly 5 bullet points explaining why Nexus Group is the right choice "
        "for a PKI/digital identity project. Each bullet must be one short, punchy sentence. "
        "Cover: deep PKI expertise, regulatory compliance track record, proven deployments, "
        "local support, and long-term partnership approach. "
        f"Tailor to the {industry} industry. "
        "Wrap in <why_nexus> tags, one bullet per line."
    )

    raw = _call_ollama(model, prompt, base_url, system=_SYS_PPT, num_predict=300)
    content = _extract_xml(raw, "why_nexus")
    bullets = [ln.lstrip("•-– ").strip() for ln in content.splitlines() if ln.strip()]
    add_bullet_slide(slide, bullets[:5])


def slide_14_references(
    prs: Presentation,
    analysis: dict[str, Any],
    config: dict[str, Any],
) -> None:
    """Slide 14: References & Case Studies. AI-generated via granite4.1:30b.

    Table: Industry / Solution / Scale / Outcome (3 rows).
    XML tag: <references>.

    Args:
        prs: The Presentation.
        analysis: RFP analysis dict.
        config: Runtime config.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "References & Case Studies")

    base_url = config["ollama_base_url"]
    model = config["ppt_model"]
    industry = analysis.get("industry_vertical", "")

    prompt = (
        "Write exactly 3 representative Nexus Group customer reference case studies. "
        "Return one pipe-delimited row per case study:\n"
        "Industry | Solution Deployed | Scale | Key Outcome (max 10 words)\n"
        f"Make at least one reference relevant to the {industry} industry. "
        "Wrap all rows in <references> tags."
    )

    raw = _call_ollama(model, prompt, base_url, system=_SYS_PPT, num_predict=300)
    content = _extract_xml(raw, "references")

    rows: list[list[str]] = []
    for line in content.splitlines():
        line = line.strip()
        if "|" in line:
            parts = [p.strip() for p in line.split("|")]
            if len(parts) >= 4:
                rows.append(parts[:4])

    if not rows:
        rows = [
            ["eGovernment", "National PKI", "20M citizens", "Nationwide digital ID launched"],
            ["Banking", "Certificate Manager", "500k certs", "PCI-DSS compliance achieved"],
            ["Telco", "Protocol Gateway", "Enterprise-wide", "Zero downtime migration"],
        ]

    add_table_slide(
        slide,
        ["Industry", "Solution", "Scale", "Key Outcome"],
        rows,
    )


def slide_15_next_steps(
    prs: Presentation,
    analysis: dict[str, Any],
    selected_products: list[dict[str, Any]],
    config: dict[str, Any],
) -> None:
    """Slide 15: Next Steps / Call to Action. AI-generated via granite4.1:30b.

    3–4 action bullets + static contact line. XML tag: <next_steps>.

    Args:
        prs: The Presentation.
        analysis: RFP analysis dict.
        selected_products: Selected product dicts.
        config: Runtime config.
    """
    slide = _new_blank_slide(prs)
    _add_teal_bar(slide)
    set_slide_title(slide, "Next Steps")

    base_url = config["ollama_base_url"]
    model = config["ppt_model"]
    slim = _slim_analysis(analysis)
    deadline = analysis.get("submission_deadline") or analysis.get("deadline") or "TBC"
    product_names = ", ".join(p.get("name", "") for p in selected_products)

    prompt = (
        "Write exactly 4 next-step action items for a presales process. "
        "Each must be a short, actionable sentence (max 12 words). "
        "Cover: (1) schedule technical deep-dive, (2) provide BoM/sizing, "
        "(3) arrange reference call, (4) submit proposal by deadline. "
        f"Submission deadline: {deadline}. Products: {product_names}. "
        "Wrap in <next_steps> tags, one action per line."
    )

    raw = _call_ollama(model, prompt, base_url, system=_SYS_PPT, num_predict=250)
    content = _extract_xml(raw, "next_steps")
    bullets = [ln.lstrip("•-– ").strip() for ln in content.splitlines() if ln.strip()]

    add_bullet_slide(slide, bullets[:4], font_size=17)

    # Static contact footer
    contact_box = slide.shapes.add_textbox(
        _BODY_LEFT, Inches(5.0), _BODY_WIDTH, Inches(0.45)
    )
    tf = contact_box.text_frame
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "Contact your Nexus Group account manager to proceed."
    _set_run_font(run, 12, True, _TEAL)


# ── Public API ─────────────────────────────────────────────────────────────────

def generate_deck(
    analysis: dict[str, Any],
    selected_products: list[dict[str, Any]],
    output_dir: str,
    config: dict[str, Any],
) -> str:
    """Generate a 15-slide branded PowerPoint deck from RFP analysis.

    Each AI slide makes exactly one Ollama call to granite4.1:30b and parses
    the response via XML markers. Static slides (1, 2, 11, 12) require no
    model calls. Total: 10 Ollama calls.

    Output saved to: output_dir/decks/RFQ{stem}_Deck_{YYYYMMDD}.pptx

    Args:
        analysis: Structured dict from analyzer.analyze_rfp(), must include
            a '_rfp_filename' key (set by main.py).
        selected_products: List of selected product dicts from product_selector.
        output_dir: Root output directory; decks/ subdirectory is created.
        config: Runtime configuration:
            - ollama_base_url (str): Ollama API base URL.
            - ppt_model (str): Model for all AI slides (default: granite4.1:30b).

    Returns:
        Absolute path to the saved .pptx file.

    Raises:
        OSError: If the output directory cannot be created or file cannot be saved.
    """
    ollama_base_url = config.get("ollama_base_url", _OLLAMA_BASE_URL)
    ppt_model = config.get("ppt_model", _PPT_MODEL)

    resolved_config: dict[str, Any] = {
        "ollama_base_url": ollama_base_url,
        "ppt_model": ppt_model,
    }

    decks_dir = Path(output_dir) / "decks"
    try:
        decks_dir.mkdir(parents=True, exist_ok=True)
    except OSError as exc:
        logger.error("Cannot create decks directory '%s': %s", decks_dir, exc)
        raise

    rfp_filename = analysis.get("_rfp_filename", "rfp")
    stem = Path(rfp_filename).stem[:20]
    date_str = datetime.now().strftime("%Y%m%d")
    filename = f"RFQ{stem}_Deck_{date_str}.pptx"
    file_path = decks_dir / filename

    logger.info("Generating PowerPoint deck: %s", file_path)

    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    # ── Slide 1: Cover (static) ──────────────────────────────────────────────
    logger.info("Slide  1/15: Cover [static]")
    slide_01_title(prs, analysis, selected_products)

    # ── Slide 2: Agenda (static) ─────────────────────────────────────────────
    logger.info("Slide  2/15: Agenda [static]")
    slide_02_agenda(prs)

    # ── Slide 3: Executive Summary [granite4.1:30b] ──────────────────────────
    logger.info("Slide  3/15: Executive Summary [%s]", ppt_model)
    slide_03_exec_summary(prs, analysis, selected_products, resolved_config)

    # ── Slide 4: Customer Challenge [granite4.1:30b] ─────────────────────────
    logger.info("Slide  4/15: Customer Challenge [%s]", ppt_model)
    slide_04_challenge(prs, analysis, resolved_config)

    # ── Slide 5: Solution Overview [granite4.1:30b] ──────────────────────────
    logger.info("Slide  5/15: Solution Overview [%s]", ppt_model)
    slide_05_solution_overview(prs, analysis, selected_products, resolved_config)

    # ── Slide 6: Product Portfolio [granite4.1:30b] ──────────────────────────
    logger.info("Slide  6/15: Product Portfolio [%s]", ppt_model)
    slide_06_product_portfolio(prs, selected_products, resolved_config)

    # ── Slide 7: Product Deep Dive [granite4.1:30b] ──────────────────────────
    logger.info("Slide  7/15: Product Deep Dive [%s]", ppt_model)
    slide_07_product_detail(prs, analysis, selected_products, resolved_config)

    # ── Slide 8: Requirements Match [granite4.1:30b] ─────────────────────────
    logger.info("Slide  8/15: Requirements Match [%s]", ppt_model)
    slide_08_requirements_match(prs, analysis, selected_products, resolved_config)

    # ── Slide 9: Compliance Summary [granite4.1:30b] ─────────────────────────
    logger.info("Slide  9/15: Compliance Summary [%s]", ppt_model)
    slide_09_compliance_summary(prs, analysis, selected_products, resolved_config)

    # ── Slide 10: Implementation Approach [granite4.1:30b] ───────────────────
    logger.info("Slide 10/15: Implementation Approach [%s]", ppt_model)
    slide_10_implementation(prs, resolved_config)

    # ── Slide 11: Project Timeline (static placeholder) ──────────────────────
    logger.info("Slide 11/15: Project Timeline [static placeholder]")
    slide_11_timeline(prs)

    # ── Slide 12: Our Team (static placeholder) ──────────────────────────────
    logger.info("Slide 12/15: Our Team [static placeholder]")
    slide_12_our_team(prs)

    # ── Slide 13: Why Nexus [granite4.1:30b] ─────────────────────────────────
    logger.info("Slide 13/15: Why Nexus [%s]", ppt_model)
    slide_13_why_nexus(prs, analysis, resolved_config)

    # ── Slide 14: References & Case Studies [granite4.1:30b] ─────────────────
    logger.info("Slide 14/15: References & Case Studies [%s]", ppt_model)
    slide_14_references(prs, analysis, resolved_config)

    # ── Slide 15: Next Steps [granite4.1:30b] ────────────────────────────────
    logger.info("Slide 15/15: Next Steps [%s]", ppt_model)
    slide_15_next_steps(prs, analysis, selected_products, resolved_config)

    # ── Save ──────────────────────────────────────────────────────────────────
    try:
        prs.save(str(file_path))
    except OSError as exc:
        logger.error("Failed to save deck '%s': %s", file_path, exc)
        raise

    resolved = str(file_path.resolve())
    logger.info("PowerPoint deck saved: %s", resolved)
    return resolved
